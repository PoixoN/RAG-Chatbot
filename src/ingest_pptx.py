"""Extract plain text from PowerPoint (.pptx) files."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _shape_text(shape) -> str:
    parts: list[str] = []
    if hasattr(shape, "text") and shape.text:
        parts.append(shape.text.strip())
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        try:
            for row in shape.table.rows:
                row_cells = [cell.text.strip() for cell in row.cells if cell.text]
                if row_cells:
                    parts.append(" | ".join(row_cells))
        except (AttributeError, ValueError):
            pass
    return "\n".join(p for p in parts if p)


def extract_pptx_text(path: Path) -> str:
    prs = Presentation(str(path))
    blocks: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in _iter_shapes(slide.shapes):
            t = _shape_text(shape)
            if t:
                slide_lines.append(t)
        if slide.has_notes_slide and slide.notes_slide:
            if slide.notes_slide.notes_text_frame and slide.notes_slide.notes_text_frame.text:
                n = slide.notes_slide.notes_text_frame.text.strip()
                if n:
                    slide_lines.append(f"[Notes] {n}")
        body = "\n".join(slide_lines)
        if body.strip():
            blocks.append(f"--- Slide {idx} ---\n{body}")
    return "\n\n".join(blocks)
